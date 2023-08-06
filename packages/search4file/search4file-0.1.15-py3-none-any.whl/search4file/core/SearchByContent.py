import asyncio
import json
import os
from pathlib import Path

import fitz
import pandas as pd
from docx import Document
from poprogress import simple_progress
from pptx import Presentation
import imghdr
from win32com.client import Dispatch, DispatchEx

from search4file.lib.unhandle_file_type import unhandle_file_types


class SearchByContent():
    # 初始化输出结果
    def __init__(self):
        # 查找结果，返回格式，3部分
        self.search_result_dict = {
            'search_path': None,
            'search_content': None,
            # 初始化、找不到时：
            'search_result': {},
            'search_problem': {"为什么搜不到后缀为doc文档的内容":"python无法读取后缀为【doc】的文档，请先用doc2docx功能，把doc文档批量转为docx文档后，再进行搜索==>>>doc2docx，批量转换的教程：https://www.bilibili.com/video/BV1so4y1H7rj"}
        }
        # 各类文件的查找结果
        self.word_result_dict = {}
        self.excel_result_dict = {}
        self.ppt_result_dict = {}
        self.pdf_result_dict = {}
        self.files_result_dict = {}
        self.search_tasks = []

    # 格式化返回内容
    def create_return_result(self, search_path, search_content):
        # 查找路径，转为绝对路径
        search_path = os.path.abspath(search_path)
        # 查询路径在哪里
        self.search_result_dict['search_path'] = search_path
        # 查询内容是什么
        self.search_result_dict['search_content'] = search_content
        # 返回纯文本的结果
        if len(self.files_result_dict) > 0:
            self.search_result_dict['search_result']['files'] = self.files_result_dict
        # 返回word的结果
        if len(self.word_result_dict) > 0:
            self.search_result_dict['search_result']['word'] = self.word_result_dict
        # 返回excel的结果
        if len(self.excel_result_dict) > 0:
            self.search_result_dict['search_result']['excel'] = self.excel_result_dict
        # 返回 pdf 的结果
        if len(self.pdf_result_dict) > 0:
            self.search_result_dict['search_result']['pdf'] = self.pdf_result_dict
        # 返回 ppt 的结果
        if len(self.ppt_result_dict) > 0:
            self.search_result_dict['search_result']['ppt'] = self.ppt_result_dict
        # 转换为json格式返回
        return json.dumps(self.search_result_dict, indent=4, ensure_ascii=False)

    # 搜索内容的逻辑
    async def search_files(self, search_path, search_content):
        try:
            # for root, dirs, files in tqdm(os.walk(search_path), unit='file', desc='搜索进度'):
            for root, dirs, files in os.walk(search_path):
                # files_list = get_files(Path(search_path).absolute())
                for file_path in simple_progress(files, desc='搜索进度'):
                    file_path = os.path.join(root, file_path)
                    # print(file_path)
                    if file_path.endswith(("docx", "doc")):  # 搜索word文件
                        self.search_tasks.append(self.search_word_file(file_path, search_content))
                    elif file_path.endswith(("xlsx", "xls")):  # 搜索excel文件
                        self.search_tasks.append(self.search_excel_file(file_path, search_content))
                    elif file_path.endswith('pdf'):  # 搜索pdf文件
                        self.search_tasks.append(self.search_pdf_file(file_path, search_content))
                    elif file_path.endswith('pptx'):  # 搜索pptx文件
                        self.search_tasks.append(self.search_ppt_file(file_path, search_content))
                    elif imghdr.what(file_path):  # 图片类文件的搜索，待开发
                        pass
                    elif file_path.endswith(unhandle_file_types):  # 待开发类文件的搜索，暂时略过
                        pass
                    else:
                        self.search_tasks.append(self.search_txt_file(file_path, search_content))  # 没有任何匹配后缀，搜索纯文本文件
            await asyncio.gather(*self.search_tasks)
        except:
            pass
        return self.create_return_result(search_path, search_content)

    ####################################
    # 搜索 纯文本 文件
    ####################################
    async def search_txt_file(self, file_path, search_content) -> None:
        """
        :param search_path:
        :param content:
        :return: dict
        """
        # 利用 open() 函数读取文件，并通过 try...except... 捕获不可读的文件格式（.zip 格式）
        try:
            with open(file_path, 'r', encoding='utf-8') as pure_file:
                for line in pure_file:
                    if search_content in line:
                        # 若是文件，则将该查询到的文件所在路径插入 final_result 空列表
                        self.files_result_dict[str(len(self.files_result_dict) + 1)] = file_path
                        break
        except:
            pass

    ####################################
    # 搜索 word 文件
    ####################################
    async def search_word_file(self, file_path, search_content):
        if Path(file_path).suffix == '.docx':
            document = Document(file_path)
            all_paragraphs = document.paragraphs
            for paragraph in all_paragraphs:
                if paragraph.text.find(search_content) >= 0:
                    self.word_result_dict[str(len(self.word_result_dict) + 1)] = file_path
                    break
        else:
            pass
            # doc_app = Dispatch('Word.Application')  # 打开word应用程序
            # # word = DispatchEx('Word.Application') # 启动独立的进程
            # doc_app.Visible = 0  # 后台运行,不显示
            # doc_app.DisplayAlerts = 0  # 不警告
            # doc_file = doc_app.Documents.Open(FileName=file_path, Encoding='gbk')
            # for para in doc_file.paragraphs:
            #     if para.text.find(search_content) >= 0:
            #         self.word_result_dict[str(len(self.word_result_dict) + 1)] = file_path
            #         break

    ####################################
    # 搜索 excel 文件
    ####################################
    async def search_excel_file(self, file_path, search_content):
        df = pd.read_excel(file_path, sheet_name=None, header=None)
        for sheet in df.values():
            for row in sheet.itertuples():
                if search_content in row:
                    self.excel_result_dict[str(len(self.excel_result_dict) + 1)] = file_path
                    break

    ####################################
    # 搜索 pdf 文件
    ####################################
    async def search_pdf_file(self, file_path, search_content):
        for page in fitz.open(file_path):  # iterate the document pages
            if page.search_for(search_content):
                self.pdf_result_dict[str(len(self.pdf_result_dict) + 1)] = file_path
                break

    ####################################
    # 搜索 ppt 文件
    ####################################
    async def search_ppt_file(self, file_path, search_content):
        ppt = Presentation(file_path)
        for slide in ppt.slides:  # > .slides 得到一个列表，包含每个列表slide
            for shape in slide.shapes:  # > slide.shapes 形状
                if shape.has_text_frame:  # shape.has_text_frame 判断是否有文字
                    # text_frame = shape.text_frame  # shape.text_frame 获取文字框
                    for paragraph in shape.text_frame.paragraphs:  # text_frame.paragraphs 获取段落
                        if search_content in paragraph.text:
                            self.ppt_result_dict[str(len(self.ppt_result_dict) + 1)] = file_path
                            break
